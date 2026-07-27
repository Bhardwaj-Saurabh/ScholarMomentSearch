"""Deck (PDF/PPTX) parsing — one slide = one unit (DESIGN.md component 3). Mirrors
paper.py's role: pure parsing over a local file, no network/DB/Qdrant. Slide
numbers are the citation locator.

Text-thin slides are flagged `needs_caption` and given an image (a rendered page
for PDF decks, the largest embedded picture for PPTX decks) for the vision-LLM
captioning task in the ingest flow (component 4) to caption before embedding —
that call is an external, retryable network operation and belongs in a Prefect
task, not here, mirroring where t_transcript sits in the video pipeline.
"""
from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

import fitz
from PIL import Image

TEXT_THIN_CHARS = 40  # a slide with less extracted text than this needs captioning
_JPEG_QUALITY = 85


@dataclass
class SlideChunk:
    slide: int                 # 1-indexed — the citation locator
    text: str                  # extracted slide text (may be thin or empty)
    needs_caption: bool        # True when text alone isn't enough to embed meaningfully
    image_jpeg: bytes | None   # slide image for captioning, when needs_caption


def parse_deck(path: Path) -> list[SlideChunk]:
    """Dispatch by extension: PDF decks render one page per slide; PPTX decks
    read shapes/text (and any embedded picture) per slide."""
    suffix = Path(path).suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf_deck(path)
    if suffix == ".pptx":
        return _parse_pptx_deck(path)
    raise ValueError(f"Unsupported deck format: {suffix!r} (expected .pdf or .pptx)")


def _to_jpeg_bytes(image_bytes: bytes) -> bytes:
    """Normalize any embedded image format to JPEG, matching the PDF path."""
    im = Image.open(BytesIO(image_bytes)).convert("RGB")
    buf = BytesIO()
    im.save(buf, format="JPEG", quality=_JPEG_QUALITY)
    return buf.getvalue()


def _parse_pdf_deck(path: Path) -> list[SlideChunk]:
    doc = fitz.open(path)
    try:
        chunks: list[SlideChunk] = []
        for i in range(doc.page_count):
            page = doc[i]
            text = page.get_text().strip()
            thin = len(text) < TEXT_THIN_CHARS
            image_jpeg = (page.get_pixmap().tobytes("jpeg", jpg_quality=_JPEG_QUALITY)
                         if thin else None)
            chunks.append(SlideChunk(slide=i + 1, text=text, needs_caption=thin,
                                     image_jpeg=image_jpeg))
        return chunks
    finally:
        doc.close()


def _parse_pptx_deck(path: Path) -> list[SlideChunk]:
    from pptx import Presentation  # heavy import kept lazy
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    chunks: list[SlideChunk] = []
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        largest_blob: bytes | None = None
        largest_area = 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                area = shape.width * shape.height
                if area > largest_area:
                    largest_area = area
                    largest_blob = shape.image.blob
        text = "\n".join(texts).strip()
        thin = len(text) < TEXT_THIN_CHARS
        image_jpeg = _to_jpeg_bytes(largest_blob) if (thin and largest_blob) else None
        chunks.append(SlideChunk(slide=i + 1, text=text, needs_caption=thin,
                                 image_jpeg=image_jpeg))
    return chunks
