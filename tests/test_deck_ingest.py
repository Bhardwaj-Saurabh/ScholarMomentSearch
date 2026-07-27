"""Component 3 (DESIGN.md) — deck PDF/PPTX -> one-slide-per-unit (src/ingest/deck.py).

Pure parsing over a local file: no network, no LLM call, no DB. Image-heavy slides
are flagged `needs_caption` with an image ready for the vision-LLM captioning task
in the ingest flow (component 4) — the network call itself lives there, mirroring
where t_transcript sits in the video pipeline, not in the parser.

SLA relevance: no direct network/queue surface, so no benchmark/sla.json row applies
standalone; feeds ingestion throughput and cross-source recall@10 once wired into
components 4 and 7. Fixture decks are built on the fly (PyMuPDF / python-pptx) and
never committed — no media files in this repo (CLAUDE.md hygiene rule).
"""
from __future__ import annotations

from io import BytesIO

import fitz
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from src.ingest import deck

LONG_LINES = [
    "This slide describes the proposed method in detail,",
    "covering the architecture, training procedure, and",
    "evaluation setup used throughout the paper.",
]
SLIDE_ONE_LINES = ["First slide body text here, well over the thin-text threshold."]
SLIDE_THREE_LINES = ["Third slide body text here, well over the thin-text threshold."]


# ── PDF decks (pymupdf: one page = one slide) ───────────────────────────────

def _build_pdf_deck(path, slides: list[list[str] | None]):
    """slides: one entry per slide — a list of lines, or None for a blank slide.
    Lines are placed individually (not a single insert_text call): PyMuPDF
    doesn't wrap a single long string, it silently clips at the page edge."""
    doc = fitz.open()
    for lines in slides:
        page = doc.new_page()
        y = 72
        for line in (lines or []):
            page.insert_text((72, y), line, fontsize=14)
            y += 18
    doc.save(str(path))
    doc.close()
    return path


def test_pdf_text_heavy_slide_needs_no_caption(tmp_path):
    pdf = _build_pdf_deck(tmp_path / "deck.pdf", [LONG_LINES])
    chunks = deck.parse_deck(pdf)
    assert len(chunks) == 1
    assert chunks[0].slide == 1
    assert "proposed method in detail" in chunks[0].text
    assert chunks[0].needs_caption is False
    assert chunks[0].image_jpeg is None


def test_pdf_blank_slide_needs_caption_with_image(tmp_path):
    pdf = _build_pdf_deck(tmp_path / "deck.pdf", [None])
    chunks = deck.parse_deck(pdf)
    assert chunks[0].needs_caption is True
    assert chunks[0].text == ""
    assert chunks[0].image_jpeg  # non-empty bytes
    Image.open(BytesIO(chunks[0].image_jpeg)).verify()  # decodable JPEG


def test_pdf_slide_numbering_across_multiple_slides(tmp_path):
    pdf = _build_pdf_deck(tmp_path / "deck.pdf",
                          [SLIDE_ONE_LINES, None, SLIDE_THREE_LINES])
    chunks = deck.parse_deck(pdf)
    assert [c.slide for c in chunks] == [1, 2, 3]
    assert chunks[1].needs_caption is True    # the blank middle slide
    assert chunks[0].needs_caption is False
    assert chunks[2].needs_caption is False


# ── PPTX decks (python-pptx: one Slide = one slide) ─────────────────────────

def _build_pptx_deck(path, slides: list[dict]):
    """slides: [{"text": str | None, "picture": bool}, ...]."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for spec in slides:
        slide = prs.slides.add_slide(blank_layout)
        if spec.get("text"):
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            box.text_frame.text = spec["text"]
        if spec.get("picture"):
            img = Image.new("RGB", (40, 30), color=(200, 50, 50))
            buf = BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(buf, Inches(1), Inches(1), Inches(2), Inches(1.5))
    prs.save(str(path))
    return path


def test_pptx_text_heavy_slide_needs_no_caption(tmp_path):
    long_text = " ".join(LONG_LINES)
    pptx = _build_pptx_deck(tmp_path / "deck.pptx", [{"text": long_text}])
    chunks = deck.parse_deck(pptx)
    assert len(chunks) == 1
    assert chunks[0].slide == 1
    assert "proposed method in detail" in chunks[0].text
    assert chunks[0].needs_caption is False
    assert chunks[0].image_jpeg is None


def test_pptx_picture_only_slide_needs_caption_with_image(tmp_path):
    pptx = _build_pptx_deck(tmp_path / "deck.pptx", [{"picture": True}])
    chunks = deck.parse_deck(pptx)
    assert chunks[0].needs_caption is True
    assert chunks[0].text == ""
    assert chunks[0].image_jpeg
    im = Image.open(BytesIO(chunks[0].image_jpeg))
    im.verify()


def test_pptx_slide_numbering_across_multiple_slides(tmp_path):
    pptx = _build_pptx_deck(tmp_path / "deck.pptx", [
        {"text": "First slide has real content here, well over the thin-text threshold."},
        {"picture": True},
        {"text": "Third slide has real content here too, well over the thin-text threshold."},
    ])
    chunks = deck.parse_deck(pptx)
    assert [c.slide for c in chunks] == [1, 2, 3]
    assert chunks[1].needs_caption is True
    assert chunks[0].needs_caption is False
    assert chunks[2].needs_caption is False


def test_unsupported_extension_raises(tmp_path):
    bogus = tmp_path / "deck.key"
    bogus.write_text("not a deck")
    with pytest.raises(ValueError):
        deck.parse_deck(bogus)
